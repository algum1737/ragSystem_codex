import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

from generator.section_generator import GeneratedSection, SECTION_DISPLAY

logger = logging.getLogger(__name__)

# python-pptx 기본 레이아웃 인덱스
_LAYOUT_TITLE_SLIDE = 0   # 표지: 제목 + 부제목
_LAYOUT_TITLE_CONTENT = 1  # 일반: 제목 + 내용


class PptxWriter:
    def write(
        self,
        sections: list[GeneratedSection],
        template_path: Path | None,
        output_path: Path,
    ) -> None:
        if template_path is not None:
            self._write_from_template(sections, template_path, output_path)
        else:
            self._write_auto(sections, output_path)

    def _write_from_template(
        self,
        sections: list[GeneratedSection],
        template_path: Path,
        output_path: Path,
    ) -> None:
        prs = Presentation(str(template_path))
        section_map = {s.section_name: s.content for s in sections}
        section_keys = list(section_map.keys())

        for idx, slide in enumerate(prs.slides):
            if idx >= len(section_keys):
                break
            section_name = section_keys[idx]
            content = section_map.get(section_name, "")
            self._fill_slide(slide, section_name, content)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info("PPTX 저장 완료 (템플릿): %s", output_path)

    def _write_auto(self, sections: list[GeneratedSection], output_path: Path) -> None:
        prs = Presentation()
        # 기본 슬라이드 크기: 와이드스크린 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        layouts = prs.slide_layouts
        # 레이아웃 인덱스가 부족한 경우 안전하게 fallback
        layout_cover = layouts[min(_LAYOUT_TITLE_SLIDE, len(layouts) - 1)]
        layout_content = layouts[min(_LAYOUT_TITLE_CONTENT, len(layouts) - 1)]

        for idx, section in enumerate(sections):
            layout = layout_cover if idx == 0 else layout_content
            slide = prs.slides.add_slide(layout)
            display = SECTION_DISPLAY.get(section.section_name, section.section_name)

            placeholders = list(slide.placeholders)
            # 제목
            if placeholders:
                try:
                    placeholders[0].text = display
                    placeholders[0].text_frame.paragraphs[0].font.size = Pt(32 if idx == 0 else 24)
                    placeholders[0].text_frame.paragraphs[0].font.bold = True
                except Exception:
                    pass
            # 내용
            if len(placeholders) > 1:
                try:
                    placeholders[1].text = section.content[:3000]
                    placeholders[1].text_frame.paragraphs[0].font.size = Pt(14)
                except Exception:
                    pass

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info("PPTX 저장 완료 (자동 생성): %s", output_path)

    def _fill_slide(self, slide, section_name: str, content: str) -> None:
        display = SECTION_DISPLAY.get(section_name, section_name)
        placeholders = list(slide.placeholders)
        if not placeholders:
            return
        try:
            placeholders[0].text = display
        except Exception:
            pass
        if len(placeholders) > 1:
            try:
                placeholders[1].text = content[:2000]
            except Exception:
                pass
